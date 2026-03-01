#!/usr/bin/env python3
"""Build a family vacation presentation for September 2026."""

import os
import time
import subprocess
import requests
import json
import base64
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

PHOTO_DIR = "/Users/robpugh/Vibe/vacation-photos"
OUTPUT = "/Users/robpugh/Vibe/family-vacation-2026.pptx"
GEMINI_API_KEY = os.environ.get("GEMINI_API_KEY", "")


# ── Photo downloading ──

def get_wikimedia_url(filename):
    """Get direct image URL from Wikimedia Commons."""
    api_url = "https://commons.wikimedia.org/w/api.php"
    params = {
        "action": "query", "titles": f"File:{filename}",
        "prop": "imageinfo", "iiprop": "url", "iiurlwidth": "1280", "format": "json",
    }
    resp = requests.get(api_url, params=params, timeout=15,
                        headers={"User-Agent": "VacationPlanBot/1.0 (vacation planning)"})
    data = resp.json()
    for page in data["query"]["pages"].values():
        if "imageinfo" in page:
            info = page["imageinfo"][0]
            return info.get("thumburl", info["url"])
    return None


def download_image(url, local_name):
    """Download image using curl."""
    path = os.path.join(PHOTO_DIR, local_name)
    if os.path.exists(path) and os.path.getsize(path) > 1000:
        print(f"  Already have {local_name}")
        return path
    print(f"  Downloading {local_name}...")
    try:
        tmp_path = path + ".tmp"
        result = subprocess.run(
            ["curl", "-sL", "-o", tmp_path, "-w", "%{http_code}",
             "-H", "User-Agent: Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36",
             url], capture_output=True, text=True, timeout=30)
        status = result.stdout.strip()
        if status != "200" or not os.path.exists(tmp_path) or os.path.getsize(tmp_path) < 1000:
            print(f"    Failed: HTTP {status}")
            if os.path.exists(tmp_path): os.remove(tmp_path)
            return None
        img = Image.open(tmp_path).convert("RGB")
        max_dim = 1920
        if max(img.size) > max_dim:
            ratio = max_dim / max(img.size)
            img = img.resize((int(img.size[0] * ratio), int(img.size[1] * ratio)), Image.LANCZOS)
        img.save(path, "JPEG", quality=85)
        os.remove(tmp_path)
        return path
    except Exception as e:
        print(f"  ERROR: {e}")
        if os.path.exists(path + ".tmp"): os.remove(path + ".tmp")
        return None


def generate_gemini_image(prompt, local_name):
    """Generate image using Gemini API."""
    path = os.path.join(PHOTO_DIR, local_name)
    if os.path.exists(path) and os.path.getsize(path) > 1000:
        print(f"  Already have {local_name}")
        return path
    print(f"  Generating {local_name} via Gemini...")
    try:
        url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-flash-exp-image-generation:generateContent?key={GEMINI_API_KEY}"
        payload = {
            "contents": [{"parts": [{"text": f"Generate a photorealistic travel photograph: {prompt}. Make it look like a professional travel magazine photo, vibrant colors, high quality, landscape orientation 16:9."}]}],
            "generationConfig": {"responseModalities": ["TEXT", "IMAGE"]}
        }
        resp = requests.post(url, json=payload, timeout=60)
        data = resp.json()
        if "candidates" in data:
            for part in data["candidates"][0]["content"]["parts"]:
                if "inlineData" in part:
                    img_data = base64.b64decode(part["inlineData"]["data"])
                    img = Image.open(BytesIO(img_data)).convert("RGB")
                    max_dim = 1920
                    if max(img.size) > max_dim:
                        ratio = max_dim / max(img.size)
                        img = img.resize((int(img.size[0] * ratio), int(img.size[1] * ratio)), Image.LANCZOS)
                    img.save(path, "JPEG", quality=90)
                    return path
        print(f"    No image in response")
        return None
    except Exception as e:
        print(f"  ERROR: {e}")
        return None


# Photos: mix of Wikimedia Commons + Gemini-generated
WIKIMEDIA_PHOTOS = {
    # Pantanal - jaguar, otter, gruta are good; swap caiman for macaw
    "pantanal_jaguar": "Jaguar (Panthera onca palustris) male Three Brothers River.JPG",
    "pantanal_otter": "Giant otter (Pteronura brasiliensis) juvenile.jpg",
    "pantanal_macaw": "Hyacinth macaw (Anodorhynchus hyacinthinus) in flight.JPG",
    "bonito_gruta": "Gruta do Lago Azul (Bonito).jpg",
    # Salvador + Morro + Boipeba
    "salvador_pelourinho": "Pelourinho Salvador Bahia 2018-0601.jpg",
    "morro_beach": "1ª Praia - Morro de São Paulo, Bahia.JPG",
    "boipeba_island": "Ilha de Boipeba, Bahia, Brasil.jpg",
    "salvador_capoeira": "Capoeira in Salvador Brazil.jpg",
    # Rio - #1 and #2 are good, keep
    "rio_panorama": "Rio de Janeiro, Brazil -08.jpg",
    "paraty_colonial": "Paraty - Rio de Janeiro (22282343219).jpg",
    # Trancoso - #3 (espelho) is good, keep
    "trancoso_espelho": "Praia do Espelho.jpg",
    "trancoso_nativos": "MARCIO FILHO PRAIA DOS NATIVOS TRANCOSO BAHIA (40977238971).jpg",
}

GEMINI_PHOTOS = {
    # Rio replacements (#3 and #4)
    "paraty_schooner": "Wooden schooner boat sailing in Paraty bay with emerald green water, small tropical islands with lush vegetation in background, mountains behind, passengers on deck enjoying the boat trip, Paraty, Brazil.",
    "ilha_grande_lopes": "Lopes Mendes beach on Ilha Grande, Brazil. Long stretch of pristine white sand, crystal clear turquoise water, gentle waves, dense Atlantic Forest vegetation on both sides, no buildings visible, one of the most beautiful beaches in the world.",
    # Trancoso replacements (#1 and #4)
    "trancoso_quadrado_colorful": "The Quadrado historic square in Trancoso, Bahia, Brazil. Row of charming colorful colonial houses painted in bright blues, yellows, greens and pinks, green grass lawn in the center, small white church at the end, tropical trees, warm golden hour light.",
    "trancoso_cliff_beach": "Dramatic view of Trancoso coastline from above, tall orange and red clay cliffs with lush green vegetation on top, wide golden sand beach below with turquoise ocean, waves breaking on shore, Bahia, Brazil.",
}


def download_all_photos():
    """Download all photos."""
    os.makedirs(PHOTO_DIR, exist_ok=True)
    photo_paths = {}

    # Wikimedia
    print("── Wikimedia Commons ──")
    seen_urls = set()
    for key, filename in WIKIMEDIA_PHOTOS.items():
        url = get_wikimedia_url(filename)
        if url:
            if url in seen_urls:
                for pk, pp in photo_paths.items():
                    if pp and WIKIMEDIA_PHOTOS.get(pk) == filename:
                        photo_paths[key] = pp
                        break
                continue
            seen_urls.add(url)
            photo_paths[key] = download_image(url, f"{key}.jpg")
            time.sleep(2)
        else:
            print(f"  Not found: {filename}")
            photo_paths[key] = None

    # Gemini
    print("\n── Gemini Generated ──")
    for key, prompt in GEMINI_PHOTOS.items():
        photo_paths[key] = generate_gemini_image(prompt, f"{key}.jpg")
        time.sleep(2)

    # Pre-generated route maps (from generate_route_maps.py)
    for key in ["map_pantanal", "map_bahia", "map_rio", "map_trancoso"]:
        path = os.path.join(PHOTO_DIR, f"{key}.jpg")
        if os.path.exists(path):
            photo_paths[key] = path
            print(f"  Route map: {key}.jpg")
        else:
            print(f"  Missing route map: {key}.jpg — run generate_route_maps.py first")
            photo_paths[key] = None

    found = sum(1 for v in photo_paths.values() if v)
    total = len(photo_paths)
    print(f"\nGot {found}/{total} photos")
    return photo_paths


# ── Color scheme ──
BG_DARK = RGBColor(0x1a, 0x1a, 0x2e)
BG_MEDIUM = RGBColor(0x16, 0x21, 0x3e)
ACCENT_GOLD = RGBColor(0xff, 0xd7, 0x00)
ACCENT_TEAL = RGBColor(0x00, 0xd2, 0xd3)
TEXT_WHITE = RGBColor(0xff, 0xff, 0xff)
TEXT_LIGHT = RGBColor(0xcc, 0xcc, 0xcc)
ACCENT_CORAL = RGBColor(0xff, 0x63, 0x48)
GREEN = RGBColor(0x2e, 0xcc, 0x71)

DEST_COLORS = {
    "pantanal": RGBColor(0x2e, 0xcc, 0x71),
    "bahia": RGBColor(0x00, 0xb8, 0xd4),
    "rio": RGBColor(0xff, 0xd7, 0x00),
    "trancoso": RGBColor(0xff, 0x63, 0x48),
}

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines,
                          font_size=16, color=TEXT_WHITE, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, is_bold, text_color, size_override) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size_override or font_size)
        p.font.color.rgb = text_color or color
        p.font.bold = is_bold
        p.font.name = font_name
        p.space_after = Pt(4)
    return txBox


def add_bullet_list(slide, left, top, width, height, items,
                    font_size=16, color=TEXT_WHITE, bullet_color=None):
    """Add a bulleted list to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, text_color) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"  {text}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = text_color or color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
        # Add bullet character
        p.text = f"\u2022  {text}"
    return txBox


def add_image_safe(slide, photo_paths, key, left, top, width, height):
    path = photo_paths.get(key)
    if path and os.path.exists(path):
        try:
            slide.shapes.add_picture(path, left, top, width, height)
            return True
        except Exception as e:
            print(f"  Image error {key}: {e}")
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x55)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"[{key}]"
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_LIGHT
    p.alignment = PP_ALIGN.CENTER
    return False


# ── Slide builders ──

def build_title_slide(prs, photo_paths):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
                "FAMILY ADVENTURE 2026", font_size=52, color=ACCENT_GOLD,
                bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(3), Inches(11), Inches(0.8),
                "September 12 - 21  |  Where should we go?", font_size=28,
                color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.6),
                "4 incredible destinations  -  YOU decide!", font_size=22,
                color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)

    dest_names = ["Pantanal", "Salvador + Morro", "Rio + Paraty", "Trancoso"]
    dest_keys = ["pantanal_jaguar", "salvador_pelourinho",
                 "rio_panorama", "trancoso_espelho"]
    start_x = Inches(2.0)
    for i, (name, key) in enumerate(zip(dest_names, dest_keys)):
        x = start_x + Inches(i * 2.5)
        add_image_safe(slide, photo_paths, key, x, Inches(5.2), Inches(1.8), Inches(1.2))
        add_textbox(slide, x, Inches(6.45), Inches(1.8), Inches(0.4),
                    name, font_size=13, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)


def build_hero_slide(prs, photo_paths, title, tagline, hero_key,
                     accent_color, map_key=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    # Hero photo on right
    add_image_safe(slide, photo_paths, hero_key,
                   Inches(6.5), Inches(0.3), Inches(6.5), Inches(4.5))
    # Color accent bar
    bar = slide.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(0.15), Inches(2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()
    # Title + tagline
    add_textbox(slide, Inches(1), Inches(1), Inches(5.2), Inches(1),
                title, font_size=40, color=accent_color, bold=True)
    add_textbox(slide, Inches(1), Inches(2.2), Inches(5.2), Inches(0.6),
                f'"{tagline}"', font_size=20, color=TEXT_LIGHT)
    # Route map below title
    if map_key:
        add_image_safe(slide, photo_paths, map_key,
                       Inches(0.5), Inches(3.0), Inches(5.5), Inches(4.2))


def build_photos_slide(prs, photo_paths, title, photo_keys, captions, accent_color):
    """Photos-only slide with captions."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_MEDIUM)
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                title, font_size=28, color=accent_color, bold=True)

    num = min(len(photo_keys), 4)
    pw = Inches(2.9)
    ph = Inches(2.2)
    gap = Inches(0.25)
    total = num * pw + (num - 1) * gap
    sx = (SLIDE_WIDTH - total) // 2

    for i, (key, cap) in enumerate(zip(photo_keys[:4], captions[:4])):
        x = sx + i * (pw + gap)
        add_image_safe(slide, photo_paths, key, x, Inches(1.3), pw, ph)
        add_textbox(slide, x, Inches(3.6), pw, Inches(0.4),
                    cap, font_size=13, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)


def build_activities_slide(prs, title, activities, accent_color):
    """Bulleted activities list - clean and readable."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_MEDIUM)
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                f"{title} — Activities", font_size=28, color=accent_color, bold=True)

    # Two columns of bullets
    left_items = activities[:len(activities)//2]
    right_items = activities[len(activities)//2:]

    for col, items in enumerate([left_items, right_items]):
        x = Inches(0.8) + col * Inches(6.2)
        y = Inches(1.3)
        bullet_items = []
        for label, desc in items:
            bullet_items.append((label, accent_color))
            bullet_items.append((f"    {desc}", TEXT_LIGHT))
        add_bullet_list(slide, x, y, Inches(5.8), Inches(5.5),
                        bullet_items, font_size=17)


def build_itinerary_slide(prs, title, days, accent_color):
    """Day-by-day itinerary table.
    days = [("1", "Sat Sep 12", "Fly SCL > GRU > Cuiaba ..."), ...]
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
                f"{title} — The Plan", font_size=28, color=accent_color, bold=True)

    num_rows = len(days) + 1  # header + data
    num_cols = 3
    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.3), Inches(0.9), Inches(12.7), Inches(6.3)
    )
    table = table_shape.table
    table.columns[0].width = Inches(0.6)
    table.columns[1].width = Inches(1.5)
    table.columns[2].width = Inches(10.6)

    # Header
    for i, h in enumerate(["Day", "Date", "Plan"]):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = ACCENT_GOLD
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER if i < 2 else PP_ALIGN.LEFT
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x0d, 0x0d, 0x1a)

    # Data rows
    for r, (day_num, date_str, plan) in enumerate(days):
        row_idx = r + 1
        bg = RGBColor(0x1f, 0x1f, 0x3a) if r % 2 == 0 else BG_DARK
        for c, (val, align) in enumerate([
            (day_num, PP_ALIGN.CENTER),
            (date_str, PP_ALIGN.CENTER),
            (plan, PP_ALIGN.LEFT),
        ]):
            cell = table.cell(row_idx, c)
            cell.text = str(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = TEXT_WHITE if c < 2 else TEXT_LIGHT
                p.font.name = "Calibri"
                p.alignment = align
                if c == 0:
                    p.font.bold = True
                    p.font.color.rgb = accent_color
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg


def build_proscons_slide(prs, title, pros, cons, cost, accent_color):
    """Pros, cons, and budget on a clean slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                f"{title} — Why / Why Not", font_size=28, color=accent_color, bold=True)

    # Pros on left
    add_textbox(slide, Inches(0.8), Inches(1.2), Inches(3.5), Inches(0.5),
                "WHY GO", font_size=18, color=GREEN, bold=True)
    pro_items = [(f"+ {p}", GREEN) for p in pros]
    add_bullet_list(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5),
                    pro_items, font_size=16)

    # Cons on right
    add_textbox(slide, Inches(7), Inches(1.2), Inches(3.5), Inches(0.5),
                "WATCH OUT", font_size=18, color=ACCENT_CORAL, bold=True)
    con_items = [(f"- {c}", ACCENT_CORAL) for c in cons]
    add_bullet_list(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(4.5),
                    con_items, font_size=16)

    # Budget at bottom
    add_textbox(slide, Inches(0.5), Inches(6.7), Inches(12), Inches(0.5),
                f"Budget: {cost}", font_size=18, color=ACCENT_GOLD, bold=True,
                alignment=PP_ALIGN.CENTER)


def build_flight_slide(prs, title, accent_color, outbound_header, outbound_flights,
                       return_header, return_info, notes):
    """
    Flight itinerary slide showing real flight options.
    outbound_flights = [("Depart", "Arrive", "Duration", "Stops", "Price", "Airline"), ...]
    return_info = string describing return options
    notes = [string, ...] footnotes
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                f"{title} — Flights", font_size=28, color=accent_color, bold=True)

    # Outbound header
    add_textbox(slide, Inches(0.5), Inches(1.1), Inches(6), Inches(0.4),
                outbound_header, font_size=16, color=ACCENT_GOLD, bold=True)

    # Build outbound table
    num_rows = len(outbound_flights) + 1  # header + data
    num_cols = 6
    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.5), Inches(1.6), Inches(8), Inches(0.35 * num_rows + 0.3)
    )
    table = table_shape.table
    col_widths = [Inches(1.1), Inches(1.1), Inches(1.2), Inches(1.5), Inches(1.1), Inches(2.0)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # Table header
    headers = ["Depart", "Arrive", "Duration", "Stops", "Price (RT)", "Airline"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = ACCENT_GOLD
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x0d, 0x0d, 0x1a)

    # Table data
    for r, flight in enumerate(outbound_flights):
        row_idx = r + 1
        bg = RGBColor(0x1f, 0x1f, 0x3a) if r % 2 == 0 else BG_DARK
        for c, val in enumerate(flight):
            cell = table.cell(row_idx, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.color.rgb = TEXT_WHITE
                p.font.name = "Calibri"
                p.alignment = PP_ALIGN.CENTER
                # Highlight nonstop and best prices
                if c == 3 and "Nonstop" in str(val):
                    p.font.color.rgb = GREEN
                elif c == 3:
                    p.font.color.rgb = TEXT_LIGHT
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg

    # Return info on right side
    table_bottom = Inches(1.6) + Inches(0.35 * num_rows + 0.3) + Inches(0.3)
    add_textbox(slide, Inches(9), Inches(1.1), Inches(4), Inches(0.4),
                return_header, font_size=16, color=ACCENT_GOLD, bold=True)
    add_textbox(slide, Inches(9), Inches(1.6), Inches(4), Inches(2.5),
                return_info, font_size=14, color=TEXT_LIGHT)

    # Notes at bottom
    if notes:
        note_lines = [(n, False, TEXT_LIGHT, 12) for n in notes]
        add_multiline_textbox(slide, Inches(0.5), Inches(6.2), Inches(12), Inches(1.2), note_lines)


def build_comparison_slide(prs):
    """Clean numeric comparison with averages."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
                "HEAD TO HEAD", font_size=36, color=ACCENT_GOLD, bold=True,
                alignment=PP_ALIGN.CENTER)

    headers = ["", "Pantanal\n+ Bonito", "Salvador\n+ Morro + Boipeba",
               "Rio + Paraty\n+ Ilha Grande", "Trancoso"]
    # Scores out of 5
    rows = [
        ("Fun Activities", [5, 4, 4, 3]),
        ("Activity Variety", [4, 4, 5, 2]),
        ("Beach / Water", [3, 5, 4, 5]),
        ("Unique Factor", [5, 4, 3, 3]),
        ("Easy Travel", [2, 3, 5, 3]),
        ("Weather", [3, 4, 4, 5]),
        ("Budget", [2, 4, 4, 1]),
    ]

    # Compute averages
    num_dest = len(headers) - 1  # minus the label column
    averages = [0.0] * num_dest
    for _, scores in rows:
        for i, s in enumerate(scores):
            averages[i] += s
    averages = [round(a / len(rows), 1) for a in averages]

    # Build table
    num_rows = len(rows) + 3  # header + data + blank + average
    table_shape = slide.shapes.add_table(
        num_rows, len(headers),
        Inches(0.5), Inches(1.1), Inches(12.3), Inches(5.5)
    )
    table = table_shape.table

    # Set column widths
    table.columns[0].width = Inches(2.0)
    for i in range(1, len(headers)):
        table.columns[i].width = Inches(2.06)

    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = ACCENT_GOLD
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x0d, 0x0d, 0x1a)

    # Color for each destination column
    dest_cols = [DEST_COLORS["pantanal"], DEST_COLORS["bahia"],
                 DEST_COLORS["rio"], DEST_COLORS["trancoso"]]

    # Data rows
    for r, (label, scores) in enumerate(rows):
        row_idx = r + 1
        # Label
        cell = table.cell(row_idx, 0)
        cell.text = label
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(15)
            p.font.color.rgb = TEXT_WHITE
            p.font.bold = True
            p.font.name = "Calibri"
        bg = RGBColor(0x1f, 0x1f, 0x3a) if r % 2 == 0 else BG_DARK
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg

        # Scores
        for c, score in enumerate(scores):
            cell = table.cell(row_idx, c + 1)
            cell.text = str(score)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(22)
                p.font.bold = True
                p.font.name = "Calibri"
                p.alignment = PP_ALIGN.CENTER
                # Color based on score
                if score >= 5:
                    p.font.color.rgb = GREEN
                elif score >= 4:
                    p.font.color.rgb = TEXT_WHITE
                elif score >= 3:
                    p.font.color.rgb = TEXT_LIGHT
                else:
                    p.font.color.rgb = ACCENT_CORAL
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg

    # Blank separator row
    sep_idx = len(rows) + 1
    for c in range(len(headers)):
        cell = table.cell(sep_idx, c)
        cell.text = ""
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x0d, 0x0d, 0x1a)

    # Average row
    avg_idx = len(rows) + 2
    cell = table.cell(avg_idx, 0)
    cell.text = "AVERAGE"
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(16)
        p.font.color.rgb = ACCENT_GOLD
        p.font.bold = True
        p.font.name = "Calibri"
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0x0d, 0x0d, 0x1a)

    # Find max average for highlighting
    max_avg = max(averages)
    for c, avg in enumerate(averages):
        cell = table.cell(avg_idx, c + 1)
        cell.text = str(avg)
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER
            p.font.color.rgb = ACCENT_GOLD if avg == max_avg else TEXT_WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x0d, 0x0d, 0x1a)



def build_voting_slide(prs, photo_paths):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_textbox(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(1),
                "WHICH ADVENTURE DO YOU WANT?", font_size=44, color=ACCENT_GOLD,
                bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(1.5), Inches(12), Inches(0.6),
                "Pick your top choice!", font_size=24,
                color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)

    options = [
        ("1. PANTANAL", "Jaguars + Crystal\nRiver Snorkeling", "pantanal_jaguar", DEST_COLORS["pantanal"]),
        ("2. SALVADOR + MORRO", "Culture, Music +\nIsland Paradise", "salvador_pelourinho", DEST_COLORS["bahia"]),
        ("3. RIO + PARATY", "Iconic City +\nColonial + Island", "rio_panorama", DEST_COLORS["rio"]),
        ("4. TRANCOSO", "Bohemian Beach\nParadise", "trancoso_espelho", DEST_COLORS["trancoso"]),
    ]
    card_w, card_h, gap = Inches(2.2), Inches(1.5), Inches(0.3)
    sx = Inches(0.6)
    for i, (name, desc, pk, color) in enumerate(options):
        x = sx + i * (card_w + gap)
        add_image_safe(slide, photo_paths, pk, x, Inches(2.5), card_w, card_h)
        add_textbox(slide, x, Inches(4.2), card_w, Inches(0.5),
                    name, font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(4.7), card_w, Inches(0.8),
                    desc, font_size=11, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)



def main():
    print("=== Family Vacation Presentation Builder ===\n")
    photo_paths = download_all_photos()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── TITLE ──
    print("\nBuilding slides...")
    build_title_slide(prs, photo_paths)

    # ── PANTANAL + BONITO (4 slides: hero, photos, activities, plan) ──
    build_hero_slide(prs, photo_paths,
        "PANTANAL + BONITO",
        "Jaguars, Piranhas & Crystal Rivers",
        "pantanal_jaguar", DEST_COLORS["pantanal"],
        map_key="map_pantanal")

    build_photos_slide(prs, photo_paths, "PANTANAL + BONITO",
        ["pantanal_jaguar", "pantanal_otter", "pantanal_macaw", "bonito_gruta"],
        ["Jaguar on the riverbank", "Giant river otter", "Hyacinth macaw in flight", "Gruta do Lago Azul, Bonito"],
        DEST_COLORS["pantanal"])

    build_activities_slide(prs, "PANTANAL + BONITO",
        [
            ("Jaguar Boat Safari", "Boat tours on the Cuiaba River with 90%+ sighting rate in September"),
            ("Piranha Fishing", "Catch and release — surprisingly easy and exciting"),
            ("Night Safari", "Boat + jeep spotting caimans with glowing eyes, nocturnal animals"),
            ("Crystal River Snorkeling", "Float 2-3 hrs in Rio da Prata's crystal water with colorful fish (Bonito)"),
            ("Underground Cave", "Gruta do Lago Azul — blue underground lake, 300 steps down"),
            ("Horseback + Birds", "Ride through the wetlands, spot capybaras, toucans, hyacinth macaws"),
        ],
        DEST_COLORS["pantanal"])

    build_itinerary_slide(prs, "PANTANAL + BONITO",
        [
            ("1", "Sat Sep 12", "Fly SCL > GRU > Cuiaba, arrive ~3-4pm. Transfer to Pantanal lodge (3-4h drive). Arrive evening, night sounds of the wetland."),
            ("2", "Sun Sep 13", "Pantanal — Early morning boat safari for jaguars (peak Sept, >90% sighting rate). Afternoon horseback ride through wetlands. Night safari for owls, caimans."),
            ("3", "Mon Sep 14", "Pantanal — Piranha fishing morning. Giant river otter spotting. Sunset bird-watching: toucans, hyacinth macaws. Night caiman spotting by boat."),
            ("4", "Tue Sep 15", "Pantanal — Second jaguar boat safari (different river stretch). Canoe through flooded forest. Capybara families along the banks."),
            ("5", "Wed Sep 16", "Pantanal — Morning wildlife walk. Afternoon free at lodge (pool, hammocks). Transfer to Cuiaba, fly to Campo Grande (1h)."),
            ("6", "Thu Sep 17", "Drive Campo Grande > Bonito (3h). Afternoon: Gruta do Lago Azul — blue underground lake, 300 steps down into the cave."),
            ("7", "Fri Sep 18", "Bonito — Rio da Prata snorkeling: float 2-3 hrs in crystal-clear river with colorful fish. Current carries you — effortless, kids 6+ OK."),
            ("8", "Sat Sep 19", "Bonito — Ecopark Porto da Ilha inflatable boat tour with waterfalls. Aquario Natural snorkeling. Evening: pack up."),
            ("9", "Sun Sep 20", "Drive Bonito > Campo Grande (3h). Fly CGR > GRU > SCL. Arrive Santiago late evening."),
        ],
        DEST_COLORS["pantanal"])

    build_proscons_slide(prs, "PANTANAL + BONITO",
        [
            "PEAK wildlife season — best month of the year",
            "Jaguar sightings nearly guaranteed (>90%)",
            "Mix: boat, jeep, horseback, hiking, snorkeling",
            "Bonito snorkeling effortless (current carries you)",
            "Family-specific tours available",
        ],
        [
            "Can hit 105°F+ (tours avoid midday heat)",
            "No beach component",
            "Pantanal lodges are pricier ($200-400/night)",
            "Long travel day from Santiago (~10h)",
            "Mosquitoes (fewer in dry season)",
        ],
        "$$$-$$$$ — Premium lodges + Bonito activities",
        DEST_COLORS["pantanal"])

    # ── SALVADOR + MORRO DE SAO PAULO + BOIPEBA (5 slides) ──
    build_hero_slide(prs, photo_paths,
        "SALVADOR + MORRO + BOIPEBA",
        "Culture, Music & Island Paradise",
        "salvador_pelourinho", DEST_COLORS["bahia"],
        map_key="map_bahia")

    build_photos_slide(prs, photo_paths, "SALVADOR + MORRO + BOIPEBA",
        ["salvador_pelourinho", "salvador_capoeira", "morro_beach", "boipeba_island"],
        ["Pelourinho, Salvador", "Capoeira in the streets", "Morro de São Paulo beach", "Boipeba island"],
        DEST_COLORS["bahia"])

    build_activities_slide(prs, "SALVADOR + MORRO + BOIPEBA",
        [
            ("Pelourinho", "UNESCO colonial center — colorful buildings, cobblestone streets, live music everywhere"),
            ("Capoeira & Culture", "Watch capoeira circles in the streets, Afro-Brazilian museums, Olodum drumming"),
            ("Bahian Food", "Acarajé, moqueca, tapioca — some of Brazil's best street food and restaurants"),
            ("Morro de São Paulo", "Car-free island, numbered beaches (2nd for party, 3rd & 4th for families), snorkeling"),
            ("Boipeba Island", "Quieter sister island — natural tide pools, mangrove boat trips, zero cars"),
            ("Boat Trips", "Island hopping between Morro and Boipeba, Gamboa beach, Garapuá natural pools"),
        ],
        DEST_COLORS["bahia"])

    build_itinerary_slide(prs, "SALVADOR + MORRO + BOIPEBA",
        [
            ("1", "Sat Sep 12", "Fly SCL > GRU > Salvador (~9-10h with connection). Arrive evening, check in at Pelourinho area."),
            ("2", "Sun Sep 13", "Salvador — Pelourinho historic center, Mercado Modelo, capoeira shows, Afro-Brazilian food tour. Live music in the evening."),
            ("3", "Mon Sep 14", "Salvador — Barra lighthouse, Farol da Barra beach, Olodum drumming. Afternoon: Ribeira neighborhood, Sorvetão ice cream."),
            ("4", "Tue Sep 15", "Catamaran Salvador > Morro de São Paulo (2.5h boat). Arrive, explore the village. Afternoon: 2nd beach swimming, sunset from the fort."),
            ("5", "Wed Sep 16", "Morro — 3rd and 4th beaches (family-friendly, calmer water). Snorkeling at natural pools. Walk to Gamboa beach at low tide."),
            ("6", "Thu Sep 17", "Morro — Full-day boat trip: Garapuá natural pools (crystal-clear swimming), Boipeba island preview, lunch on the water."),
            ("7", "Fri Sep 18", "Boat Morro > Boipeba (~40min). Quieter island, gorgeous beaches. Bike or walk to Cueira and Tassimirim beaches. Tide pool exploration."),
            ("8", "Sat Sep 19", "Boipeba — Mangrove boat trip, Ponta dos Castelhanos (remote beach accessible only by boat). Last swim. Evening: stargazing."),
            ("9", "Sun Sep 20", "Boat Boipeba > Valença, transfer to Salvador airport. Fly SSA > GRU > SCL. Arrive Santiago late evening."),
        ],
        DEST_COLORS["bahia"])

    build_proscons_slide(prs, "SALVADOR + MORRO + BOIPEBA",
        [
            "Rich Afro-Brazilian culture — music, dance, food",
            "Two car-free islands back to back",
            "Morro beaches are family-perfect (calm, shallow)",
            "Boipeba is off the beaten path, uncrowded",
            "September = dry season, warm (27°C), great weather",
            "Bahian food is incredible",
        ],
        [
            "Salvador > Morro is 2.5h boat (can be rough seas)",
            "Salvador needs safety awareness (like Rio)",
            "Less variety than adventure options",
            "Connection flight from Santiago (~9-10h)",
            "Morro can feel touristy at 2nd beach",
        ],
        "$$-$$$ — Mid-range, islands are affordable",
        DEST_COLORS["bahia"])

    # ── RIO + PARATY + ILHA GRANDE (5 slides) ──
    build_hero_slide(prs, photo_paths,
        "RIO + PARATY + ILHA GRANDE",
        "Iconic City, Colonial Charm & Island Paradise",
        "rio_panorama", DEST_COLORS["rio"],
        map_key="map_rio")

    build_photos_slide(prs, photo_paths, "RIO + PARATY + ILHA GRANDE",
        ["rio_panorama", "paraty_colonial", "paraty_schooner", "ilha_grande_lopes"],
        ["Rio de Janeiro panorama", "Paraty colonial streets", "Schooner boat trip, Paraty", "Lopes Mendes beach, Ilha Grande"],
        DEST_COLORS["rio"])

    build_activities_slide(prs, "RIO + PARATY + ILHA GRANDE",
        [
            ("Christ the Redeemer", "Corcovado Mountain, Sugarloaf cable car, Tijuca Forest hikes"),
            ("Rio Beaches", "Copacabana, Ipanema — iconic boardwalks, acai bowls, beach culture"),
            ("Paraty Colonial Town", "UNESCO cobblestone streets, colorful houses, cachaca distilleries"),
            ("Island Hopping", "Schooner boat trip from Paraty — 65 islands, emerald water, snorkeling"),
            ("Ilha Grande Beaches", "Lopes Mendes (top 10 world beach), Lagoa Azul for snorkeling"),
            ("Waterfall Hikes", "Forest trails to hidden waterfalls on the car-free island"),
        ],
        DEST_COLORS["rio"])

    build_itinerary_slide(prs, "RIO + PARATY + ILHA GRANDE",
        [
            ("1", "Sat Sep 12", "Fly SCL > Rio NONSTOP (4h, arrive 10:30am!). Check in Copacabana/Ipanema. Afternoon: beach walk, acai, settle in."),
            ("2", "Sun Sep 13", "Rio — Corcovado (Christ the Redeemer) morning. Sugarloaf cable car afternoon. Dinner in Lapa neighborhood."),
            ("3", "Mon Sep 14", "Rio — Tijuca Forest hike to waterfalls. Ipanema beach afternoon. Optional: bike ride along the coastline."),
            ("4", "Tue Sep 15", "Drive Rio > Paraty (5h scenic Costa Verde coastal road). Arrive afternoon, explore the UNESCO colonial cobblestone streets."),
            ("5", "Wed Sep 16", "Paraty — Full-day schooner boat trip: 65 islands, emerald water, snorkeling stops, beach lunch. Kids can jump off the boat."),
            ("6", "Thu Sep 17", "Paraty — Morning kayak tour through mangroves. Afternoon: cachaca distillery visit, free time in colonial center. Evening: waterfront dinner."),
            ("7", "Fri Sep 18", "Ferry Paraty > Ilha Grande (~1h). Car-free island! Settle in Abraao village. Afternoon: easy trail to Praia Preta beach."),
            ("8", "Sat Sep 19", "Ilha Grande — Boat trip to Lagoa Azul (snorkeling in blue lagoon) + Lopes Mendes beach (top 10 in the world). Full day on the water."),
            ("9", "Sun Sep 20", "Ilha Grande — Morning waterfall hike through Atlantic Forest. Ferry back to Mangaratiba, drive to GIG. Fly home nonstop (~4h)."),
        ],
        DEST_COLORS["rio"])

    build_proscons_slide(prs, "RIO + PARATY + ILHA GRANDE",
        [
            "EASIEST logistics — direct 4h nonstop flight!",
            "World-famous landmarks (Christ, Sugarloaf)",
            "Great variety: city + colonial + island",
            "Paraty boat trips are magical (65 islands)",
            "Ilha Grande is car-free paradise",
            "September = dry season, pleasant temps",
        ],
        [
            "Rio needs safety awareness (tourist areas OK)",
            "5-hour drive Rio > Paraty",
            "Ilha Grande is rustic (basic infrastructure)",
            "Less exotic/unique than other options",
        ],
        "$$-$$$ — Wide range, all price points",
        DEST_COLORS["rio"])

    # ── TRANCOSO (5 slides) ──
    build_hero_slide(prs, photo_paths,
        "TRANCOSO",
        "Brazil's Most Beautiful Beach Village",
        "trancoso_espelho", DEST_COLORS["trancoso"],
        map_key="map_trancoso")

    build_photos_slide(prs, photo_paths, "TRANCOSO",
        ["trancoso_quadrado_colorful", "trancoso_nativos", "trancoso_espelho", "trancoso_cliff_beach"],
        ["The Quadrado historic square", "Praia dos Nativos", "Praia do Espelho", "Trancoso cliff coastline"],
        DEST_COLORS["trancoso"])

    build_activities_slide(prs, "TRANCOSO",
        [
            ("The Quadrado", "Historic colorful square — restaurants, shops, green lawn, small white church"),
            ("Praia do Espelho", "One of Brazil's top 5 beaches — natural tide pools, coral reefs"),
            ("Horseback Riding", "Ride along the coast and through Atlantic Forest trails"),
            ("Beach Buggy Tours", "Discover hidden coves and beaches along the coast"),
            ("Film Festival", "Trancoso Film Festival in September — outdoor screenings under the stars"),
            ("Kayaking + Snorkeling", "Paddle up the Trancoso River, snorkel at offshore coral reefs"),
        ],
        DEST_COLORS["trancoso"])

    build_itinerary_slide(prs, "TRANCOSO",
        [
            ("1", "Sat Sep 12", "Fly SCL > GRU > Porto Seguro (~8-9h). Drive 1.5h to Trancoso. Arrive evening, dinner at the Quadrado."),
            ("2", "Sun Sep 13", "Trancoso — Explore the Quadrado (colorful historic square, white church). Beach day at Praia dos Nativos. Kids play in the square."),
            ("3", "Mon Sep 14", "Day trip to Praia do Espelho — one of Brazil's top 5 beaches. Natural tide pools at low tide, coral reef snorkeling."),
            ("4", "Tue Sep 15", "Horseback riding along the coast and through Atlantic Forest trails. Afternoon: beach buggy tour to hidden coves."),
            ("5", "Wed Sep 16", "Bike tour through forest trails to secluded beaches. Afternoon: Trancoso Film Festival outdoor screening under the stars."),
            ("6", "Thu Sep 17", "Kayak up the Trancoso River through mangroves. Afternoon: snorkeling at offshore coral reefs. Evening: Quadrado restaurants."),
            ("7", "Fri Sep 18", "Day trip to Arraial d'Ajuda — charming village nearby, different beaches, water park option for kids. Return for sunset."),
            ("8", "Sat Sep 19", "Visit Pataxo indigenous village (learn about their culture, crafts). Last beach afternoon. Farewell dinner at the Quadrado."),
            ("9", "Sun Sep 20", "Drive Trancoso > Porto Seguro (1.5h). Fly BPS > GRU > SCL. Arrive Santiago late evening."),
        ],
        DEST_COLORS["trancoso"])

    build_proscons_slide(prs, "TRANCOSO",
        [
            "Stunningly beautiful beaches",
            "The Quadrado is magical",
            "Very safe, relaxed village",
            "Perfect weather in September",
            "Amazing restaurant scene",
        ],
        [
            "Most expensive option (Brazil's St. Tropez)",
            "Can feel slow after 4-5 days for active kids",
            "Limited variety beyond beach + dining",
            "Connections from Santiago (~8-9h)",
            "More adult-oriented than adventure",
        ],
        "$$$$ — Premium, shoulder season helps",
        DEST_COLORS["trancoso"])

    # ── COMPARISON ──
    build_comparison_slide(prs)

    # ── VOTING ──
    build_voting_slide(prs, photo_paths)

    prs.save(OUTPUT)
    print(f"\nSaved: {OUTPUT}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
